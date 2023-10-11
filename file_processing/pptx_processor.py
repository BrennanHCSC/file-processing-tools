from file_processor_strategy import FileProcessorStrategy
from pptx import Presentation


class PptxFileProcessor(FileProcessorStrategy):
    def __init__(self, file_path: str) -> None:
        super().__init__(file_path)
        self.metadata = {}

    def process(self) -> None:
        ppt = Presentation(self.file_path)
        self.metadata.update({'text': self.extract_text_from_pptx(ppt)})
        self.metadata.update({'author': ppt.core_properties.author})
        self.metadata.update({'last_modified_by': ppt.core_properties.last_modified_by})
        self.metadata.update({"num_slides": len(ppt.slides)})

        # Other core properties to include: https://python-pptx.readthedocs.io/en/latest/api/presentation.html#coreproperties-objects
        # keywords, language, subject, version

    def save(self, output_path: str = None) -> None:
        ppt = Presentation(self.file_path)

        # Update the core properties (metadata)
        cp = ppt.core_properties
        cp.author = self.metadata.get('author', cp.author)
        cp.last_modified_by = self.metadata.get('last_modified_by', cp.last_modified_by)
        
        save_path = output_path or self.file_path
        ppt.save(save_path)


    def save_as_pdf(self, output_path: str = None) -> None:
        import comtypes.client
        import os
        import sys
        input = sys.argv[1]
        output = sys.argv[2]

        input = os.path.abspath(self.file_path)
        output = os.path.abspath(output_path)
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")

        slides = powerpoint.Presentations.Open(input)

        slides.SaveAs(output, 32)

        slides.Close()

    @staticmethod
    def extract_text_from_pptx(ppt: Presentation) -> str:
        try:
            full_text = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)
                    if shape.has_table:
                        for r in shape.table.rows:
                            s = ""
                            for c in r.cells:
                                s += c.text_frame.text + " | "
                            full_text.append(s)
            return '\n'.join(full_text)
        except Exception as e:
            print(f"Error encountered while opening or processing {file_path}: {e}")
            return None
