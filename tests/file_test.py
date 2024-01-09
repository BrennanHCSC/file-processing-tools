import sys, os
sys.path.append(os.path.join(sys.path[0],'file_processing'))

def test_docx_text():
    from file_processing.file import File
    docx_1 = File('tests/resources/test_files/HealthCanadaOverviewFromWikipedia.docx')
    docx_2 = File('tests/resources/test_files/SampleReport.docx')
    assert len(docx_1.metadata['text']) == 1631
    assert len(docx_2.metadata['text']) == 3220


def test_docx_author():
    from file_processing.file import File
    from docx import Document

    # Paths of test files
    test_docx_1_path = 'tests/resources/test_files/HealthCanadaOverviewFromWikipedia.docx'
    test_docx_2_path = 'tests/resources/test_files/SampleReport.docx'

    # Save paths of output files
    test_docx_1_save_path = 'tests/resources/test_files/HealthCanadaOverviewFromWikipediaModified.docx'
    test_docx_2_save_path = 'tests/resources/test_files/SampleReportModified.docx'

    # Arbitrary test author names
    test_author_1 = 'Test Author One'
    test_author_2 = 'Second Test Author'

    # Update test docx files to have test names
    for x in [(test_docx_1_path, test_docx_1_save_path, test_author_1), 
              (test_docx_2_path, test_docx_2_save_path, test_author_2)]:
        doc = Document(x[0])
        doc.core_properties.author = x[2]
        doc.save(x[1])

    # Test author names match
    docx_1 = File(test_docx_1_save_path)
    docx_2 = File(test_docx_2_save_path)
    assert docx_1.metadata['author'] == test_author_1
    assert docx_2.metadata['author'] == test_author_2

    os.remove(test_docx_1_save_path)
    os.remove(test_docx_2_save_path)


def test_docx_last_modified_by():
    from file_processing.file import File
    from docx import Document

    # Paths of test files
    test_docx_1_path = 'tests/resources/test_files/HealthCanadaOverviewFromWikipedia.docx'
    test_docx_2_path = 'tests/resources/test_files/SampleReport.docx'

    # Save paths of output files
    test_docx_1_save_path = 'tests/resources/test_files/HealthCanadaOverviewFromWikipediaModified.docx'
    test_docx_2_save_path = 'tests/resources/test_files/SampleReportModified.docx'

    # Arbitrary test last_modified_by names
    test_last_modified_by_1 = 'Test last_modified_by One'
    test_last_modified_by_2 = 'last_modified_by Test Author'

    # Update test docx files to have test names
    for x in [(test_docx_1_path, test_docx_1_save_path, test_last_modified_by_1), 
              (test_docx_2_path, test_docx_2_save_path, test_last_modified_by_2)]:
        doc = Document(x[0])
        doc.core_properties.last_modified_by = x[2]
        doc.save(x[1])

    # Test last_modified_by names match
    docx_1 = File(test_docx_1_save_path)
    docx_2 = File(test_docx_2_save_path)
    assert docx_1.metadata['last_modified_by'] == test_last_modified_by_1
    assert docx_2.metadata['last_modified_by'] == test_last_modified_by_2

    os.remove(test_docx_1_save_path)
    os.remove(test_docx_2_save_path)


def test_docx_locked():
    from file_processing.file import File

    docx_1 = File('tests/resources/test_files/SampleReport_Locked.docx')
    assert docx_1.metadata["has_password"] == True

    docx_2 = File('tests/resources/test_files/HealthCanadaOverviewFromWikipedia_Locked.docx')
    assert docx_2.metadata["has_password"] == True


def test_save_docx_metadata():
    from file_processing.file import File
    from docx import Document
    
    test_docx_path = 'tests/resources/test_files/HealthCanadaOverviewFromWikipedia.docx'
    copy_test_docx_path = 'tests/resources/test_files/HealthCanadaOverviewFromWikipedia_copy.docx'
    
    # Copying file
    with open(test_docx_path, 'rb') as src_file:
        with open(copy_test_docx_path, 'wb') as dest_file:
            dest_file.write(src_file.read())

    try:
        # Arbitrary test author and last_modified_by names
        test_author = 'New Author'
        test_last_modified_by = 'Modified Author'
        
        # Load and change metadata via File object
        docx_file = File(copy_test_docx_path)
        docx_file.metadata['author'] = test_author
        docx_file.metadata['last_modified_by'] = test_last_modified_by
        
        # Save changes
        docx_file.save()
        
        # Load document again to check if the changes were saved correctly
        doc = Document(copy_test_docx_path)
        
        # Assert if changes are correctly reflected
        assert doc.core_properties.author == test_author
        assert doc.core_properties.last_modified_by == test_last_modified_by

    finally:
        # Clean up by removing the copied file after the test is done
        os.remove(copy_test_docx_path)

def test_docx_invalid_save_location():
    import pytest
    from file_processing.file import File
    from errors import FileProcessingFailedError
    
    docx_file = File('tests/resources/test_files/HealthCanadaOverviewFromWikipedia.docx')
    with pytest.raises(FileProcessingFailedError):
        docx_file.processor.save('/non_existent_folder/HealthCanadaOverviewFromWikipedia.docx')

def test_docx_corrupted_file_processing():
    import pytest
    from file_processing.file import File
    from errors import FileCorruptionError
    with pytest.raises(FileCorruptionError) as exc_info:
        File("tests/resources/test_files/HealthCanadaOverviewFromWikipedia_corrupted.docx")


def test_pdf_ocr_text_found():
    from file_processing.file import File
    pdf_1 = File('tests/resources/test_files/SampleReportScreenShot.pdf', use_ocr=True)
    ocr_text = pdf_1.metadata['ocr_text']
    assert len(ocr_text) > 0


def test_pdf_locked():
    from file_processing.file import File

    pdf_1 = File('tests/resources/test_files/SampleReport_Locked.pdf')
    assert pdf_1.metadata["has_password"] == True

    pdf_2 = File('tests/resources/test_files/ArtificialNeuralNetworksForBeginners_Locked.pdf')
    assert pdf_2.metadata["has_password"] == True

def test_pdf_invalid_save_location():
    import pytest
    from file_processing.file import File
    from errors import FileProcessingFailedError
    
    pdf_file = File('tests/resources/test_files/ArtificialNeuralNetworksForBeginners_Locked.pdf')
    with pytest.raises(FileProcessingFailedError):
        pdf_file.processor.save('/non_existent_folder/ArtificialNeuralNetworksForBeginners_Locked.pdf')

def test_pdf_corrupted_file_processing():
    import pytest
    from file_processing.file import File
    from errors import FileProcessingFailedError
    with pytest.raises(FileProcessingFailedError) as exc_info:
        File("tests/resources/test_files/SampleReportScreenShot_corrupted.pdf")


def test_msg_text():
    from file_processing.file import File
    msg = File('tests/resources/test_files/Test Email.msg')
    msg_text = msg.metadata['text']
    assert msg_text == 'Body text.\r\n\r\n \r\n\r\n'

    
def test_msg_subject():
    from file_processing.file import File
    msg = File('tests/resources/test_files/Test Email.msg')
    msg_subject = msg.metadata['subject']
    assert msg_subject == 'Test Email'

    
def test_msg_date():
    from file_processing.file import File
    msg = File('tests/resources/test_files/Test Email.msg')
    msg_date = msg.metadata['date']
    assert msg_date == 'Mon, 18 Sep 2023 13:57:16 -0400'

    
def test_msg_sender():
    from file_processing.file import File
    msg = File('tests/resources/test_files/Test Email.msg')
    msg_sender = msg.metadata['sender']
    assert msg_sender == '"Burnett, Taylen (HC/SC)" <Taylen.Burnett@hc-sc.gc.ca>'

def test_save_msg_metadata():
    from file_processing.file import File

    test_msg_path = 'tests/resources/test_files/Test Email.msg'
    copy_test_msg_path = 'tests/resources/test_files/Test Email_copy.msg'
    
    # Copying file
    with open(test_msg_path, 'rb') as src_file:
        with open(copy_test_msg_path, 'wb') as dest_file:
            dest_file.write(src_file.read())

    try:
        # Load via File object
        msg_file = File(copy_test_msg_path)
        
        # Save changes
        msg_file.save()

        # Load document again to check if the changes were saved correctly
        msg = File(copy_test_msg_path)
        
        # Assert if file correctly saved
        assert msg.metadata['sender'] == '"Burnett, Taylen (HC/SC)" <Taylen.Burnett@hc-sc.gc.ca>'
        assert msg.metadata['date'] == 'Mon, 18 Sep 2023 13:57:16 -0400'
        assert msg.metadata['subject'] == 'Test Email'
        assert msg.metadata['text'] == 'Body text.\r\n\r\n \r\n\r\n'

    finally:
        os.remove(copy_test_msg_path)

def test_msg_invalid_save_location():
    import pytest
    from file_processing.file import File
    from errors import FileProcessingFailedError
    
    msg_file = File('tests/resources/test_files/Test Email.msg')
    with pytest.raises(FileProcessingFailedError):
        msg_file.processor.save('/non_existent_folder/Test Email.msg')

def test_msg_corrupted_file_processing():
    import pytest
    from file_processing.file import File
    from errors import FileProcessingFailedError
    with pytest.raises(FileProcessingFailedError) as exc_info:
        File("tests/resources/test_files/Test Email_corrupted.msg")


def test_png_format():
    from file_processing.file import File
    png_1 = File('tests/resources/test_files/Health_Canada_logo.png')
    png_2 = File('tests/resources/test_files/MapCanada.png')
    assert png_1.metadata['original_format'] == 'GIF'
    assert png_2.metadata['original_format'] == 'PNG'

    
def test_png_mode():
    from file_processing.file import File
    png_1 = File('tests/resources/test_files/Health_Canada_logo.png')
    png_2 = File('tests/resources/test_files/MapCanada.png')
    assert png_1.metadata['mode'] == 'P'
    assert png_2.metadata['mode'] == 'RGBA'

    
def test_png_width():
    from file_processing.file import File
    png_1 = File('tests/resources/test_files/Health_Canada_logo.png')
    png_2 = File('tests/resources/test_files/MapCanada.png')
    assert png_1.metadata['width'] == 303
    assert png_2.metadata['width'] == 3000

    
def test_png_height():
    from file_processing.file import File
    png_1 = File('tests/resources/test_files/Health_Canada_logo.png')
    png_2 = File('tests/resources/test_files/MapCanada.png')
    assert png_1.metadata['height'] == 40
    assert png_2.metadata['height'] == 2408

def test_save_png_metadata():
    from file_processing.file import File
    
    test_jpeg_path = 'tests/resources/test_files/Health_Canada_logo.png'
    copy_test_jpeg_path = 'tests/resources/test_files/Health_Canada_logo_copy.png'
    
    # Copying file
    with open(test_jpeg_path, 'rb') as src_file:
        with open(copy_test_jpeg_path, 'wb') as dest_file:
            dest_file.write(src_file.read())

    try:
        # Load via File object
        jpeg_file = File(copy_test_jpeg_path)
        
        # Save changes
        jpeg_file.save()
        
        # Load document again to check if the changes were saved correctly
        jpeg = File(copy_test_jpeg_path)
        
        # Assert if file correctly saved
        assert jpeg.metadata['height'] == 40
        assert jpeg.metadata['width'] == 303
        assert jpeg.metadata['mode'] == 'P'

    finally:
        # Clean up by removing the copied file after the test is done
        os.remove(copy_test_jpeg_path)

def test_png_invalid_save_location():
    import pytest
    from file_processing.file import File
    from errors import FileProcessingFailedError
    
    png_file = File('tests/resources/test_files/Health_Canada_logo.png')
    with pytest.raises(FileProcessingFailedError):
        png_file.processor.save('/non_existent_folder/Health_Canada_logo.png')

def test_png_corrupted_file_processing():
    import pytest
    from file_processing.file import File
    from errors import FileProcessingFailedError
    with pytest.raises(FileProcessingFailedError) as exc_info:
        File("tests/resources/test_files/MapCanada_corrupted.png")
    
    
def test_xlsx_sheets():
    from file_processing.file import File 
    exceldoc = File('tests/resources/test_files/Test_excel_file.xlsx')
    exceldoc_sheetnames = exceldoc.metadata['sheet_names']
    assert exceldoc_sheetnames == ['Sheet1', 'Sheet2', 'Sheet3']

    
def test_xlsx_activesheet():
    from file_processing.file import File
    exceldoc = File('tests/resources/test_files/Test_excel_file.xlsx')
    exceldoc_activesheet = exceldoc.metadata['active_sheet']
    assert str(exceldoc_activesheet) == "Sheet3"

def test_xlsx_data():
    from file_processing.file import File
    exceldoc = File('tests/resources/test_files/Test_excel_file.xlsx')
    assert len(exceldoc.metadata['data']['Sheet1']) == 10
    assert len(exceldoc.metadata['data']['Sheet2']) == 11
    assert len(exceldoc.metadata['data']['Sheet3']) == 21

def test_xlsx_last_modified_by():
    from file_processing.file import File
    exceldoc = File('tests/resources/test_files/Test_excel_file.xlsx')
    assert exceldoc.metadata['last_modified_by'] == 'Burnett, Taylen (HC/SC)'

def test_xlsx_creator():
    from file_processing.file import File
    exceldoc = File('tests/resources/test_files/Test_excel_file.xlsx')
    assert exceldoc.metadata['creator'] == 'Burnett, Taylen (HC/SC)'

def test_xlsx_locked():
    from file_processing.file import File

    exceldoc = File('tests/resources/test_files/Test_excel_file_Locked.xlsx')
    assert exceldoc.metadata["has_password"] == True

    exceldoc_2 = File('tests/resources/test_files/StructureofCanadianFederalGovFromWikipedia_Locked.xlsx')
    assert exceldoc_2.metadata["has_password"] == True

def test_save_xlsx_metadata():
    from file_processing.file import File
    from openpyxl import load_workbook
    
    test_exc_path = 'tests/resources/test_files/Test_excel_file.xlsx'
    copy_test_exc_path = 'tests/resources/test_files/Test_excel_file_copy.xlsx'
    
    # Copying file
    with open(test_exc_path, 'rb') as src_file:
        with open(copy_test_exc_path, 'wb') as dest_file:
            dest_file.write(src_file.read())

    try:
        # Arbitrary test author and last_modified_by names
        test_creator = 'New Creator'
        test_last_modified_by = 'Modified Creator'
        
        # Load and change metadata via File object
        exc_file = File(copy_test_exc_path)
        exc_file.metadata['creator'] = test_creator
        exc_file.metadata['last_modified_by'] = test_last_modified_by
        
        # Save changes
        exc_file.save()
        
        # Load document again to check if the changes were saved correctly
        exceldoc = load_workbook(copy_test_exc_path)
        
        # Assert if changes are correctly reflected
        assert exceldoc.properties.creator == test_creator
        assert exceldoc.properties.lastModifiedBy == test_last_modified_by

    finally:
        # Clean up by removing the copied file after the test is done
        os.remove(copy_test_exc_path)

def test_xlsx_invalid_save_location():
    import pytest
    from file_processing.file import File
    from errors import FileProcessingFailedError
    
    xlsx_file = File('tests/resources/test_files/StructureofCanadianFederalGovFromWikipedia_Locked.xlsx')
    with pytest.raises(FileProcessingFailedError):
        xlsx_file.processor.save('/non_existent_folder/StructureofCanadianFederalGovFromWikipedia_Locked.xlsx')

def test_xlsx_corrupted_file_processing():
    import pytest
    from file_processing.file import File
    from errors import FileCorruptionError
    with pytest.raises(FileCorruptionError) as exc_info:
        File("tests/resources/test_files/Test_excel_file_corrupted.xlsx")

    
def test_pptx_text():
    from file_processing.file import File
    pptx_1 = File('tests/resources/test_files/HealthCanadaOverviewFromWikipedia.pptx')
    pptx_2 = File('tests/resources/test_files/SampleReport.pptx')
    assert len(pptx_1.metadata['text']) == 1655
    assert len(pptx_2.metadata['text']) == 2037


def test_pptx_author():
    from file_processing.file import File
    from pptx import Presentation

    # Paths of test files
    test_pptx_1_path = 'tests/resources/test_files/HealthCanadaOverviewFromWikipedia.pptx'
    test_pptx_2_path = 'tests/resources/test_files/SampleReport.pptx'

    # Save paths of test files
    test_pptx_1_save_path = 'tests/resources/test_files/HealthCanadaOverviewFromWikipediaModified.pptx'
    test_pptx_2_save_path = 'tests/resources/test_files/SampleReportModified.pptx'

    # Arbitrary test author names
    test_author_1 = 'Test Author One'
    test_author_2 = 'Second Test Author'

    # Update test pptx files to have test names
    for x in [(test_pptx_1_path, test_pptx_1_save_path, test_author_1), 
              (test_pptx_2_path, test_pptx_2_save_path, test_author_2)]:
        ppt = Presentation(x[0])
        ppt.core_properties.author = x[2]
        ppt.save(x[1])

    # Test author names match
    pptx_1 = File(test_pptx_1_save_path)
    pptx_2 = File(test_pptx_2_save_path)
    assert pptx_1.metadata['author'] == test_author_1
    assert pptx_2.metadata['author'] == test_author_2

    os.remove(test_pptx_1_save_path)
    os.remove(test_pptx_2_save_path)


def test_pptx_last_modified_by():
    from file_processing.file import File
    from pptx import Presentation

    # Paths of test files
    test_pptx_1_path = 'tests/resources/test_files/HealthCanadaOverviewFromWikipedia.pptx'
    test_pptx_2_path = 'tests/resources/test_files/SampleReport.pptx'
    
    # Save paths of test files
    test_pptx_1_save_path = 'tests/resources/test_files/HealthCanadaOverviewFromWikipediaModified.pptx'
    test_pptx_2_save_path = 'tests/resources/test_files/SampleReportModified.pptx'

    # Arbitrary test last_modified_by names
    test_last_modified_by_1 = 'Test last_modified_by One'
    test_last_modified_by_2 = 'last_modified_by Test Author'

    # Update test pptx files to have test names
    for x in [(test_pptx_1_path, test_pptx_1_save_path, test_last_modified_by_1), 
              (test_pptx_2_path, test_pptx_2_save_path, test_last_modified_by_2)]:
        ppt = Presentation(x[0])
        ppt.core_properties.last_modified_by = x[2]
        ppt.save(x[1])

    # Test last_modified_by names match
    pptx_1 = File(test_pptx_1_save_path)
    pptx_2 = File(test_pptx_2_save_path)
    assert pptx_1.metadata['last_modified_by'] == test_last_modified_by_1
    assert pptx_2.metadata['last_modified_by'] == test_last_modified_by_2

    os.remove(test_pptx_1_save_path)
    os.remove(test_pptx_2_save_path)

    
def test_pptx_num_slides():
    from file_processing.file import File
    pptx_1 = File('tests/resources/test_files/HealthCanadaOverviewFromWikipedia.pptx')
    pptx_2 = File('tests/resources/test_files/SampleReport.pptx')
    assert pptx_1.metadata['num_slides'] == 4
    assert pptx_2.metadata['num_slides'] == 5


def test_pptx_locked():
    from file_processing.file import File

    pptx_1 = File('tests/resources/test_files/SampleReport_Locked.pptx')
    assert pptx_1.metadata["has_password"] == True

    pptx_2 = File('tests/resources/test_files/HealthCanadaOverviewFromWikipedia_Locked.pptx')
    assert pptx_2.metadata["has_password"] == True
    

def test_save_ppt_metadata():
    from file_processing.file import File
    from pptx import Presentation
    test_ppt_path = 'tests/resources/test_files/HealthCanadaOverviewFromWikipedia.pptx'
    copy_test_ppt_path = 'tests/resources/test_files/HealthCanadaOverviewFromWikipedia_copy.pptx'
    
    # Copying file
    with open(test_ppt_path, 'rb') as src_file:
        with open(copy_test_ppt_path, 'wb') as dest_file:
            dest_file.write(src_file.read())

    try:
        # Arbitrary test author and last_modified_by names
        test_author = 'New Author'
        test_last_modified_by = 'Modified Author'
        
        # Load and change metadata via File object
        ppt_file = File(copy_test_ppt_path)
        ppt_file.metadata['author'] = test_author
        ppt_file.metadata['last_modified_by'] = test_last_modified_by
        
        # Save changes
        ppt_file.save()
        
        # Load document again to check if the changes were saved correctly
        ppt = Presentation(copy_test_ppt_path)
        
        # Assert if changes are correctly reflected
        assert ppt.core_properties.author == test_author
        assert ppt.core_properties.last_modified_by == test_last_modified_by

    finally:
        # Clean up by removing the copied file after the test is done
        os.remove(copy_test_ppt_path)

def test_pptx_invalid_save_location():
    import pytest
    from file_processing.file import File
    from errors import FileProcessingFailedError
    
    pptx_file = File('tests/resources/test_files/HealthCanadaOverviewFromWikipedia_Locked.pptx')
    with pytest.raises(FileProcessingFailedError):
        pptx_file.processor.save('/non_existent_folder/HealthCanadaOverviewFromWikipedia_Locked.pptx')

def test_pptx_corrupted_file_processing():
    import pytest
    from file_processing.file import File
    from errors import FileCorruptionError
    with pytest.raises(FileCorruptionError) as exc_info:
        File("tests/resources/test_files/HealthCanadaOverviewFromWikipedia_corrupted.pptx")


def test_rtf_text():
    from file_processing.file import File
    rtf_1 = File('tests/resources/test_files/Test_for_RTF.rtf')
    assert len(rtf_1.metadata['text']) == 5306


def test_save_rtf_metadata():
    from file_processing.file import File
    
    test_rtf_path = 'tests/resources/test_files/Test_for_RTF.rtf'
    copy_test_rtf_path = 'tests/resources/test_files/Test_for_RTF_copy.rtf'
    
    # Copying file
    with open(test_rtf_path, 'rb') as src_file:
        with open(copy_test_rtf_path, 'wb') as dest_file:
            dest_file.write(src_file.read())

    try:
        
        # Load via File object
        rtf_file = File(copy_test_rtf_path)
        
        # Save
        rtf_file.save()
        
        # Assert if .txt correctly saved
        assert len(rtf_file.metadata['text']) == 5306

    finally:
        # Clean up by removing the copied file after the test is done
        os.remove(copy_test_rtf_path)


def test_rtf_invalid_save_location():
    import pytest
    from file_processing.file import File
    from errors import FileProcessingFailedError
    
    rtf_file = File('tests/resources/test_files/Test_for_RTF.rtf')
    with pytest.raises(FileProcessingFailedError):
        rtf_file.processor.save('/non_existent_folder/Test_for_RTF.rtf')

def test_rtf_corrupted_file_processing():
    import pytest
    from file_processing.file import File
    from errors import FileProcessingFailedError
    with pytest.raises(FileProcessingFailedError) as exc_info:
        File("tests/resources/test_files/Test_for_RTF_corrupted.rtf")

                 
def test_html_text():
    from file_processing.file import File
    txt_1 = File('tests/resources/test_files/Health - Canada.ca.html')
    assert len(txt_1.metadata['text']) == 165405

def test_html_num_lines():
    # indirectly tests lines attribute
    from file_processing.file import File
    txt_1 = File('tests/resources/test_files/Health - Canada.ca.html')
    assert txt_1.metadata['num_lines'] == 3439
    
def test_html_num_words():
    # indirectly tests words attribute
    from file_processing.file import File
    txt_1 = File('tests/resources/test_files/Health - Canada.ca.html')
    assert txt_1.metadata['num_words'] == 11162

def test_save_html_metadata():
    from file_processing.file import File
    
    test_html_path = 'tests/resources/test_files/Health - Canada.ca.html'
    copy_test_html_path = 'tests/resources/test_files/Health - Canada.ca_copy.html'
    
    # Copying file
    with open(test_html_path, 'rb') as src_file:
        with open(copy_test_html_path, 'wb') as dest_file:
            dest_file.write(src_file.read())

    try:
        
        # Load via File object
        html_file = File(copy_test_html_path)
        
        # Save
        html_file.save()
        
        # Assert if .txt correctly saved
        assert len(html_file.metadata['text']) == 165405

    finally:
        # Clean up by removing the copied file after the test is done
        os.remove(copy_test_html_path)

def test_html_invalid_save_location():
    import pytest
    from file_processing.file import File
    from errors import FileProcessingFailedError
    
    html_file = File('tests/resources/test_files/Health - Canada.ca.html')
    with pytest.raises(FileProcessingFailedError):
        html_file.processor.save('/non_existent_folder/Health - Canada.ca.html')

def test_html_corrupted_file_processing():
    import pytest
    from file_processing.file import File
    from errors import FileProcessingFailedError
    with pytest.raises(FileProcessingFailedError) as exc_info:
        File("tests/resources/test_files/Health - Canada.ca_corrupted.html")

    
def test_xml_text():
    from file_processing.file import File
    txt_1 = File('tests/resources/test_files/Sample.xml')
    assert len(txt_1.metadata['text']) == 4429

def test_xml_num_lines():
    # indirectly tests lines attribute
    from file_processing.file import File
    txt_1 = File('tests/resources/test_files/Sample.xml')
    assert txt_1.metadata['num_lines'] == 120
    
def test_xml_num_words():
    # indirectly tests words attribute
    from file_processing.file import File
    txt_1 = File('tests/resources/test_files/Sample.xml')
    assert txt_1.metadata['num_words'] == 336

def test_save_xml_metadata():
    from file_processing.file import File
    
    test_xml_path = 'tests/resources/test_files/Sample.xml'
    copy_test_xml_path = 'tests/resources/test_files/Sample_copy.xml'
    
    # Copying file
    with open(test_xml_path, 'rb') as src_file:
        with open(copy_test_xml_path, 'wb') as dest_file:
            dest_file.write(src_file.read())

    try:
        
        # Load via File object
        xml_file = File(copy_test_xml_path)
        
        # Save
        xml_file.save()
        
        # Assert if .txt correctly saved
        assert len(xml_file.metadata['text']) == 4429

    finally:
        # Clean up by removing the copied file after the test is done
        os.remove(copy_test_xml_path)

def test_xml_invalid_save_location():
    import pytest
    from file_processing.file import File
    from errors import FileProcessingFailedError
    
    xml_file = File('tests/resources/test_files/Sample.xml')
    with pytest.raises(FileProcessingFailedError):
        xml_file.processor.save('/non_existent_folder/Sample.xml')

def test_xml_corrupted_file_processing():
    import pytest
    from file_processing.file import File
    from errors import FileProcessingFailedError
    with pytest.raises(FileProcessingFailedError) as exc_info:
        File("tests/resources/test_files/Sample_corrupted.xml")


def test_jpeg_format():
    from file_processing.file import File
    jpeg_1 = File('tests/resources/test_files/HealthCanada.jpeg')
    jpeg_2 = File('tests/resources/test_files/MapCanada.jpg')
    assert jpeg_1.metadata['original_format'] == 'JPEG'
    assert jpeg_2.metadata['original_format'] == 'JPEG'
    
def test_jpeg_mode():
    from file_processing.file import File
    jpeg_1 = File('tests/resources/test_files/HealthCanada.jpeg')
    jpeg_2 = File('tests/resources/test_files/MapCanada.jpg')
    assert jpeg_1.metadata['mode'] == 'RGB'
    assert jpeg_2.metadata['mode'] == 'RGB'
    
def test_jpeg_width():
    from file_processing.file import File
    jpeg_1 = File('tests/resources/test_files/HealthCanada.jpeg')
    jpeg_2 = File('tests/resources/test_files/MapCanada.jpg')
    assert jpeg_1.metadata['width'] == 474
    assert jpeg_2.metadata['width'] == 4489

def test_jpeg_height():
    from file_processing.file import File
    jpeg_1 = File('tests/resources/test_files/HealthCanada.jpeg')
    jpeg_2 = File('tests/resources/test_files/MapCanada.jpg')
    assert jpeg_1.metadata['height'] == 262
    assert jpeg_2.metadata['height'] == 2896

def test_save_jpeg_metadata():
    from file_processing.file import File
    
    test_jpeg_path = 'tests/resources/test_files/HealthCanada.jpeg'
    copy_test_jpeg_path = 'tests/resources/test_files/HealthCanada_copy.jpeg'
    
    # Copying file
    with open(test_jpeg_path, 'rb') as src_file:
        with open(copy_test_jpeg_path, 'wb') as dest_file:
            dest_file.write(src_file.read())

    try:
        # Load via File object
        jpeg_file = File(copy_test_jpeg_path)
        
        # Save changes
        jpeg_file.save()
        
        # Load document again to check if the changes were saved correctly
        jpeg = File(copy_test_jpeg_path)
        
        # Assert if file correctly saved
        assert jpeg.metadata['height'] == 262
        assert jpeg.metadata['width'] == 474
        assert jpeg.metadata['mode'] == 'RGB'

    finally:
        # Clean up by removing the copied file after the test is done
        os.remove(copy_test_jpeg_path)

def test_jpeg_invalid_save_location():
    import pytest
    from file_processing.file import File
    from errors import FileProcessingFailedError
    
    jpeg_file = File('tests/resources/test_files/HealthCanada.jpeg')
    with pytest.raises(FileProcessingFailedError):
        jpeg_file.processor.save('/non_existent_folder/HealthCanada.jpeg')

def test_jpeg_corrupted_file_processing():
    import pytest
    from file_processing.file import File
    from errors import FileProcessingFailedError
    with pytest.raises(FileProcessingFailedError) as exc_info:
        File("tests/resources/test_files/MapCanada_corrupted.jpg")

        
def test_csv_text():
    from file_processing.file import File
    csv_1 = File('tests/resources/test_files/2021_Census_English.csv')
    csv_2 = File('tests/resources/test_files/Approved_Schools_2023_10_01.csv')
    assert len(csv_1.metadata['text']) == 4463562
    assert len(csv_2.metadata['text']) == 1274028

def test_csv_num_rows():
    from file_processing.file import File
    csv_1 = File('tests/resources/test_files/2021_Census_English.csv')
    csv_2 = File('tests/resources/test_files/Approved_Schools_2023_10_01.csv')
    assert csv_1.metadata['num_rows'] == 36835
    assert csv_2.metadata['num_rows'] == 5385

def test_csv_num_cols():
    from file_processing.file import File
    csv_1 = File('tests/resources/test_files/2021_Census_English.csv')
    csv_2 = File('tests/resources/test_files/Approved_Schools_2023_10_01.csv')
    assert csv_1.metadata['num_cols'] == 23
    assert csv_2.metadata['num_cols'] == 13

def test_csv_num_cells():
    from file_processing.file import File
    csv_1 = File('tests/resources/test_files/2021_Census_English.csv')
    csv_2 = File('tests/resources/test_files/Approved_Schools_2023_10_01.csv')
    assert csv_1.metadata['num_cells'] == 847205
    assert csv_2.metadata['num_cells'] == 70005

def test_csv_empty_cells():
    from file_processing.file import File
    csv_1 = File('tests/resources/test_files/2021_Census_English.csv')
    csv_2 = File('tests/resources/test_files/Approved_Schools_2023_10_01.csv')
    assert csv_1.metadata['empty_cells'] == 253932
    assert csv_2.metadata['empty_cells'] == 0

def test_save_csv_metadata():
    from file_processing.file import File
    
    test_csv_path = 'tests/resources/test_files/2021_Census_English.csv'
    copy_test_csv_path = 'tests/resources/test_files/2021_Census_English_copy.csv'
    
    # Copying file
    with open(test_csv_path, 'rb') as src_file:
        with open(copy_test_csv_path, 'wb') as dest_file:
            dest_file.write(src_file.read())

    try:
        
        # Load via File object
        csv_file = File(copy_test_csv_path)
        
        # Save
        csv_file.save()
        
        # Assert if changes correctly saved
        assert len(csv_file.metadata['text']) == 4463562

    finally:
        # Clean up by removing the copied file after the test is done
        os.remove(copy_test_csv_path)

def test_csv_invalid_save_location():
    import pytest
    from file_processing.file import File
    from errors import FileProcessingFailedError
    csv_file = File('tests/resources/test_files/2021_Census_English.csv')
    with pytest.raises(FileProcessingFailedError):
        csv_file.processor.save('/non_existent_folder/2021_Census_English.csv')

def test_csv_corrupted_file_processing():
    import pytest
    from file_processing.file import File
    from errors import FileProcessingFailedError
    with pytest.raises(FileProcessingFailedError) as exc_info:
        File("tests/resources/test_files/2021_Census_English_corrupted.csv")
        

def test_zip_num_files():
    from file_processing.file import File
    zip_1 = File('tests/resources/test_files/SampleReport.zip')
    zip_2 = File('tests/resources/test_files/Empty.zip')
    assert zip_1.metadata['num_files'] == 3
    assert zip_2.metadata['num_files'] == 0

def test_zip_file_types():
    from file_processing.file import File
    zip_1 = File('tests/resources/test_files/SampleReport.zip')
    zip_2 = File('tests/resources/test_files/Empty.zip')
    assert zip_1.metadata['file_types'] == {'docx': 2, 'pptx': 1}
    assert zip_2.metadata['file_types'] == {}


def test_zip_file_names():
    from file_processing.file import File
    zip_1 = File('tests/resources/test_files/SampleReport.zip')
    zip_2 = File('tests/resources/test_files/Empty.zip')
    assert zip_1.metadata['file_names'] == ['SampleReport.docx', 'SampleReport.pptx', 'HealthCanadaOverviewFromWikipedia.docx']
    assert zip_2.metadata['file_names'] == []

def test_zip_extraction():
    from file_processing.file import File
    import shutil
    zip_file = File('tests/resources/test_files/SampleReport.zip')
    zip_file.processor.extract()

    extraction_dir = 'tests/resources/test_files/SampleReport'
    assert os.path.isdir(extraction_dir)

    extracted_files = os.listdir(extraction_dir)
    expected_files = ['SampleReport.docx', 'SampleReport.pptx', 'HealthCanadaOverviewFromWikipedia.docx']
    assert set(extracted_files) == set(expected_files)

    shutil.rmtree(extraction_dir)

def test_zip_save():
    from file_processing.file import File
    import tempfile
    import zipfile
    with tempfile.TemporaryDirectory() as temp_dir:
        original_zip_path = 'tests/resources/test_files/SampleReport.zip'
        saved_zip_path = os.path.join(temp_dir, 'SavedSampleReport.zip')

        zip_file = File(original_zip_path)

        zip_file.processor.save(saved_zip_path)

        assert os.path.exists(saved_zip_path)

        with zipfile.ZipFile(original_zip_path, 'r') as original_zip, zipfile.ZipFile(saved_zip_path, 'r') as saved_zip:
            assert set(original_zip.namelist()) == set(saved_zip.namelist()) # Check contents are still the same

def test_zip_invalid_save_location():
    import pytest
    from file_processing.file import File
    from errors import FileProcessingFailedError
    zip_file = File('tests/resources/test_files/SampleReport.zip')
    with pytest.raises(FileProcessingFailedError):
        zip_file.processor.save('/non_existent_folder/SavedSampleReport.zip')

def test_zip_corrupted_file_processing():
    import pytest
    from file_processing.file import File
    from errors import FileProcessingFailedError
    with pytest.raises(FileProcessingFailedError) as exc_info:
        File("tests/resources/test_files/SampleReport_corrupted.zip")

def test_json_num_keys():
    from file_processing.file import File
    json_1 = File('tests/resources/test_files/coffee.json')
    json_2 = File('tests/resources/test_files/sample.json')
    assert json_1.metadata['num_keys'] == 15
    assert json_2.metadata['num_keys'] == 9

def test_json_key_names():
    from file_processing.file import File
    json_1 = File('tests/resources/test_files/coffee.json')
    json_2 = File('tests/resources/test_files/sample.json')
    assert json_1.metadata['key_names'] == ['quiz', 'sport', 'q1', 'question', 'options', 'answer', 'maths', 'q1', 'question', 'options', 'answer', 'q2', 'question', 'options', 'answer']
    assert json_2.metadata['key_names'] == ['array', 'boolean', 'color', 'null', 'number', 'object', 'a', 'c', 'string']

def test_json_empty_values():
    from file_processing.file import File
    json_1 = File('tests/resources/test_files/coffee.json')
    json_2 = File('tests/resources/test_files/sample.json')
    assert json_1.metadata['empty_values'] == 0
    assert json_2.metadata['empty_values'] == 1

def test_save_json_metadata():
    from file_processing.file import File

    test_json_path = 'tests/resources/test_files/coffee.json'
    copy_test_json_path = 'tests/resources/test_files/coffee_copy.json'

    # Copying file
    with open(test_json_path, 'rb') as src_file:
        with open(copy_test_json_path, 'wb') as dest_file:
            dest_file.write(src_file.read())

    try:

        # Load via File object
        json_file = File(copy_test_json_path)
        
        # Save
        json_file.save()
        
        # Assert if changes correctly saved
        assert json_file.metadata['num_keys'] == 15

    finally:
        # Clean up by removing the copied file after the test is done
        os.remove(copy_test_json_path)

def test_json_invalid_save_location():
    import pytest
    from file_processing.file import File
    from errors import FileProcessingFailedError
    json_file = File('tests/resources/test_files/coffee.json')
    with pytest.raises(FileProcessingFailedError):
        json_file.processor.save('/non_existent_folder/coffee.json')

def test_json_corrupted_file_processing():
    import pytest
    from file_processing.file import File
    from errors import FileCorruptionError
    with pytest.raises(FileCorruptionError) as exc_info:
        File("tests/resources/test_files/coffee_corrupted.json")